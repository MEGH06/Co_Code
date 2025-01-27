import os
from pptx import Presentation
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.image as mpimg

def extract_text_from_pptx(file_path):
    """Extract text from a PowerPoint file."""
    prs = Presentation(file_path)
    text = []
    
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = " ".join(run.text for run in paragraph.runs)
                    slide_text.append(paragraph_text)
        if slide_text:
            text.append("\n".join(slide_text))
    
    return text


def extract_tables_from_pptx(file_path):
    """Extract tables from a PowerPoint file as DataFrames."""
    prs = Presentation(file_path)
    tables = []
    
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                df = pd.DataFrame(table_data[1:], columns=table_data[0])  # Convert to DataFrame with headers
                tables.append({"slide_number": slide_number, "table_data": df})
    
    return tables


def extract_images_from_pptx(file_path, output_dir="extracted_images"):
    """Extract images from a PowerPoint file and save them to a directory."""
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(file_path)
    images = []
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type == 13:  # Image type
                image = shape.image
                image_bytes = image.blob
                ext = image.ext
                image_path = f"{output_dir}/slide_{slide_idx}_image_{shape_idx}.{ext}"
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                images.append({"slide_number": slide_idx, "image_path": image_path})
    
    return images


def preview_extracted_images(images):
    """Preview extracted images using matplotlib."""
    for image_info in images:
        img_path = image_info["image_path"]
        img = mpimg.imread(img_path)
        plt.imshow(img)
        plt.title(f"Slide {image_info['slide_number']}: {os.path.basename(img_path)}")
        plt.axis('off')
        plt.show()


def extract_pptx_content(file_path):
    """Extract all content (text, tables, images) from a PowerPoint file."""
    content = {"text": [], "tables": [], "images": []}
    
    # Extract text
    content["text"] = extract_text_from_pptx(file_path)
    
    # Extract tables
    tables = extract_tables_from_pptx(file_path)
    for table_info in tables:
        content["tables"].append({
            "slide_number": table_info["slide_number"],
            "table_data": table_info["table_data"]
        })
    
    # Extract images
    images = extract_images_from_pptx(file_path)
    content["images"] = images
    
    return content


def print_tables(tables):
    """Print tables in a readable format."""
    for i, table_info in enumerate(tables):
        try:
            print(f"\nSlide {table_info['slide_number']} - Table {i+1}:")
            print(table_info["table_data"].to_string(index=False))
        except Exception as e:
            print(f"Error processing table {i+1} on Slide {table_info['slide_number']}: {e}")
            continue

# Example usage
file_path = r"{path_to_pptx_file}"
content = extract_pptx_content(file_path)




